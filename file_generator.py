import os
import random
from io import BytesIO
from docx import Document
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import RGBColor as docx_RGBColor 
from openpyxl import Workbook
from openpyxl.drawing.image import Image as xlsImage
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from PyPDF2 import PdfWriter, PdfReader
from PyPDF2.generic import NameObject, ArrayObject, DictionaryObject, TextStringObject
from PIL import Image, ImageDraw, ImageFont
import tempfile
from app import app
from app import logging
from pptx.util import Inches
from pptx.dml.color import RGBColor
from app.db import *
import subprocess
import shutil

# File storage folder
FILE_STORAGE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'data', 'generated_files')

if not os.path.exists(FILE_STORAGE):
    os.makedirs(FILE_STORAGE)

# Load malicious URLs from a file
def load_malicious_urls():
    try:
        with open("app/data/malicious_urls.txt", "r") as file:
            return [line.strip() for line in file if line.strip()]
    except FileNotFoundError:
        return []

# Generate a random URL (malicious or clean based on user's selection)
def generate_random_url(url_type="malicious"):
    if url_type == "malicious":
        malicious_urls = load_malicious_urls()
        if malicious_urls:
            return random.choice(malicious_urls)
        else:
            return "http://gmai.com"  # Fallback malicious URL
    elif url_type == "clean":
        # Generate a clean random URL from a predefined set
        domains = ["example.com", "cleanwebsite.org", "safeurl.net"]
        return f"https://{random.randint(1, 100)}.{random.choice(domains)}"
    else:
        return None  # No URL

# Generate all file types
def generate_all_files(url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission):
    file_types = ['pdf', 'docx', 'pptx','xlsx', 'exe', 'rtf', 'jpg', 'png', 'bmp', 'gif', 'tiff']
    for file_type in file_types:
        generate_file(file_type, url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission)

# Generate a file based on type and selected features (URLs, Images, Scripts, Video, Audio, Link, 3D)
def generate_file(file_type, url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission):
    logging.info(f"Values received by generate_file: {file_type, url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission}")
    
    if file_type == 'pdf':
        logging.info("Calling the pdf generator")
        filename = generate_pdf(url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission)
        

        #logging.info(url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission)
    elif file_type == 'docx':
        filename = generate_docx(url_type, include_image, include_sensitive_link)
    elif file_type == 'pptx':
        filename = generate_pptx(url_type, include_image, include_sensitive_link)
    elif file_type == 'rtf':
        filename = generate_rtf(url_type, include_image, include_sensitive_link)
    elif file_type == 'xlsx':
        filename = generate_xlsx(url_type, include_image, include_sensitive_link)
    elif file_type in ['jpg', 'png', 'bmp', 'gif', 'tiff']:
        filename = generate_image(file_type, url_type, include_sensitive_link)

    elif file_type == 'exe':
        filename = generate_exe(url_type)
    
    elif file_type == 'dylib':
        filename = generate_dylib(url_type)

    elif file_type == 'elf':
        filename = generate_elf(url_type)


    # After generating the file, save its details to the database
    logging.info(f"file name before savnig to database {filename}")
    save_generated_file_to_db(
        filename, file_type, url_type, include_image, include_sensitive_link, include_script, include_video, include_audio, include_3d, include_pdf, include_external_app, include_data_submission
    )

    return filename

#Generate PDF
def generate_pdf(url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission):
    logging.info(f"The variables as seen inside generate_pdf: {url_type, include_image, include_script, include_video, include_audio, include_sensitive_link, include_3d, include_pdf, include_external_app, include_data_submission}")
    
    filename = f"generated_pdf_file_{random.randint(1000, 9999)}.pdf"
    file_path = os.path.join(FILE_STORAGE, filename)
    
    # Create a buffer to hold the PDF data in memory
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer)
    
    logging.info(f"file path: {file_path}")
    
    # Basic PDF content
    pdf.drawString(100, 750, "PDF with Advanced Features")

    # Include URL if specified
    if url_type != 'none':
        random_url = generate_random_url(url_type)
        if random_url is None:
            random_url = "No URL provided"
        pdf.drawString(80, 320, f"URL: {random_url}")
        pdf.linkURL(random_url, (80, 310, 380, 330), relative=0)

    # Include sensitive link if specified
    if include_sensitive_link == "on":
        sensitive_link = r"\\server\share"
        if sensitive_link is None:
            sensitive_link = "No link provided"
        pdf.drawString(80, 580, f"Server Share: {sensitive_link}")
        pdf.linkURL(sensitive_link, (80, 570, 380, 590), relative=0)

    # Embed image if selected
    if include_image == "on":
        # Generate image path
        image_name = generate_image("jpg", url_type, include_sensitive_link)
        random_image_path = os.path.join(FILE_STORAGE, image_name)

        # Add the image to the PDF
        add_image_to_pdf(pdf, random_image_path, x=380, y=500, width=180, height=180)
        
    # Embed JavaScript code if specified
    if include_script == "on":
        script = "app.alert('This is an embedded JavaScript test!');"
        writer = PdfWriter()
        writer.add_js(script)
        pdf.setFont("Helvetica", 10)
        pdf.drawString(80, 450, "Embedded JavaScript:")
        pdf.drawString(80, 430, script)

    # Save the initial PDF without annotations
    pdf.save()

    # Now create the PDF file object for writing
    buffer.seek(0)
    pdf_reader = PdfReader(buffer)
    pdf_writer = PdfWriter()

    # Create a page object from the PDF
    page = pdf_reader.pages[0]
    
    # Add annotations for video, audio, 3D, etc.
    annotations = []

    # Link to local video
    if include_video == "on":
        video_file_path = "file:///home/kali/Desktop/cp_demo_server/app/data/pdf_assets/sample_mp4_file.mp4"
        video_annotation = DictionaryObject({
            NameObject("/Subtype"): NameObject("/Link"),
            NameObject("/A"): DictionaryObject({
                NameObject("/S"): NameObject("/URI"),
                NameObject("/URI"): TextStringObject(video_file_path)
            })
        })
        annotations.append(video_annotation)

    # Link to local audio
    if include_audio == "on":
        audio_file_path = "file:///home/kali/Desktop/cp_demo_server/app/data/pdf_assets/sample_mp3_file.mp3"
        audio_annotation = DictionaryObject({
            NameObject("/Subtype"): NameObject("/Link"),
            NameObject("/A"): DictionaryObject({
                NameObject("/S"): NameObject("/URI"),
                NameObject("/URI"): TextStringObject(audio_file_path)
            })
        })
        annotations.append(audio_annotation)

    # Embed 3D Art (Dummy Data)
    if include_3d == "on":
        u3d_data = b"dummy 3D model data for testing purposes"
        u3d_annotation = DictionaryObject({
            NameObject("/Subtype"): NameObject("/3D"),
            NameObject("/3DName"): TextStringObject("Dummy 3D Model"),
            NameObject("/3DData"): ArrayObject([TextStringObject(u3d_data)]),
        })
        annotations.append(u3d_annotation)

    # Add external PDF if selected
    if include_pdf == "on":
        external_pdf_path = "file:///home/kali/Desktop/cp_demo_server/app/data/pdf_assets/sample_pdf_file.pdf"
        pdf_annotation = DictionaryObject()
        pdf_annotation.update({
            NameObject("/Subtype"): NameObject("/Link"),
            NameObject("/Contents"): TextStringObject("Click to open External PDF"),
            NameObject("/A"): DictionaryObject({
                NameObject("/S"): NameObject("/URI"),
                NameObject("/URI"): TextStringObject(external_pdf_path)
            })
        })
        annotations.append(pdf_annotation)

    # Add external application launch if selected
    if include_external_app == "on":
        app_path = "file:///path/to/external/application"
        app_annotation = DictionaryObject()
        app_annotation.update({
            NameObject("/Subtype"): NameObject("/Link"),
            NameObject("/Contents"): TextStringObject("Launch External Application"),
            NameObject("/A"): DictionaryObject({
                NameObject("/S"): NameObject("/URI"),
                NameObject("/URI"): TextStringObject(app_path)
            })
        })
        annotations.append(app_annotation)

    # Add data submission if selected
    if include_data_submission == "on":
        submission_url = "http://example.com/submit_data"
        submission_annotation = DictionaryObject()
        submission_annotation.update({
            NameObject("/Subtype"): NameObject("/Widget"),
            NameObject("/T"): TextStringObject("Submit Data"),
            NameObject("/A"): DictionaryObject({
                NameObject("/S"): NameObject("/URI"),
                NameObject("/URI"): TextStringObject(submission_url)
            })
        })
        annotations.append(submission_annotation)

    # Now, add the annotations to the page (only if there are annotations)
    if annotations:
        page[NameObject("/Annots")] = ArrayObject(annotations)

    # Write the page with annotations back into the PDF
    pdf_writer.add_page(page)

    # Write to the output file using a proper file object
    with open(file_path, "wb") as f:
        pdf_writer.write(f)

    buffer.close()

    return filename  # Return the filename of the generated PDF

def add_image_to_pdf(pdf, image_path, x, y, width, height):
    """
    Add an image to a PDF using ReportLab.

    Args:
        pdf (canvas.Canvas): The ReportLab canvas instance.
        image_path (str): Path to the image file.
        x (int): X-coordinate for the image.
        y (int): Y-coordinate for the image.
        width (int): Width of the image in the PDF.
        height (int): Height of the image in the PDF.
    """
    try:
        # Open the image and save it as a temporary PNG file
        img = Image.open(image_path)
        img_byte_arr = BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            img.save(tmp_file.name, format='PNG')  # Save as PNG format
            tmp_file_path = tmp_file.name

        # Embed the image into the PDF using the temporary file path
        pdf.drawImage(tmp_file_path, x, y, width, height)

        # Clean up the temporary file
        os.remove(tmp_file_path)
    except Exception as e:
        logging.info(f"Error adding image to PDF: {e}")

# Generate DOCX
def generate_docx(url_type, include_image, include_sensitive_link):
    """
    Generates a DOCX file based on user input.
    
    Args:
        url_type (str): Type of URL to include (e.g., "safe", "malicious").
        include_image (bool): Whether to include an image in the document.
        include_sensitive_link (bool): Whether to include a sensitive link.

    Returns:
        str: The filename of the generated DOCX file.
    """
    filename = f"generated_docx_file_{random.randint(1000, 9999)}.docx"
    file_path = os.path.join(FILE_STORAGE, filename)

    # Debugging logs
    logging.debug(f"Generating DOCX with parameters - URL Type: {url_type}, Include Image: {include_image}, Include Sensitive Link: {include_sensitive_link}")

    # Create a new Word document
    doc = Document()
    doc.add_heading("Generated Word Document", level=1)
    doc.add_paragraph("This document might include active contents.")

    if url_type != "none":
        # Add a random URL as a clickable link
        random_url = generate_random_url()  # Assuming this function exists and returns a valid URL
        para = doc.add_paragraph("Visit this link: ")
        create_hyperlink(para, random_url)

    # Conditionally include sensitive link
    if include_sensitive_link == "on":
        sensitive_link = "\\\\server\\share"
        para = doc.add_paragraph("Sensitive link: ")
        create_hyperlink(para, sensitive_link)
        logging.debug("Included sensitive link in the document.")

    # Conditionally include an image
    if include_image == "on":
        image_name = generate_image("jpg", url_type, include_sensitive_link)
        random_image_path = os.path.join(FILE_STORAGE, image_name)
        if os.path.exists(random_image_path):
            doc.add_picture(random_image_path)
            logging.debug("Included image in the document.")
        else:
            logging.warning(f"Image file not found: {random_image_path}")

    # Save the document
    doc.save(file_path)
    logging.info(f"DOCX file created: {filename}")

    return filename

def create_hyperlink(paragraph, url):
    # Create the hyperlink XML element
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), 'rId1')  # We will handle relationships later
    
    # Create a run with the hyperlink text
    r = OxmlElement('w:r')
    t = OxmlElement('w:t')
    t.text = url
    r.append(t)
    hyperlink.append(r)

    # Add the hyperlink to the paragraph
    paragraph._element.append(hyperlink)

    # Style the hyperlink (blue and underlined)
    run = paragraph.add_run(url)
    run.font.color.rgb = docx_RGBColor(0, 0, 255)  # Blue color
    run.font.underline = True  # Underline

    return hyperlink

# Generate PowerPoint Presentation (.pptx) with random image, and URL
def generate_pptx(url_type, include_image, include_sensitive_link):
    filename = f"generated_docx_file_{random.randint(1000, 9999)}.pptx"
    file_path = os.path.join(FILE_STORAGE, filename)

    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = "PowerPoint with Resources"

    # Add random URL if specified
    if url_type != 'none':
        random_url = generate_random_url(url_type)
        
        # Check if the generated URL is None and provide a fallback or skip
        if random_url is not None:
            # Add a textbox for the URL and make it clickable
            text_box = slide.shapes.placeholders[1]
            text_frame = text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = random_url
            
            # Adding hyperlink to the text
            run = p.runs[0]
            run.hyperlink.address = random_url
        else:
            logging.info("No URL generated, skipping.")

        
    # Include sensitive link if specified
    if include_sensitive_link == "on":
        sensitive_link = r"\\server\share"  # Replace with your actual server share link
        
        # Add a textbox for the sensitive link and make it clickable
        text_box = slide.shapes.placeholders[0]
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = sensitive_link
        
        # Adding hyperlink to the sensitive link
        run = p.runs[0]
        run.hyperlink.address = sensitive_link

    # Embed image if selected
    if include_image == "on":
        image_name = generate_image("jpg", url_type, include_sensitive_link)
        random_image_path = os.path.join(FILE_STORAGE, image_name)
        slide.shapes.add_picture(random_image_path, Inches(0.1), Inches(0.1), width=Inches(3))

    # Save the PowerPoint file
    ppt.save(file_path)

    
    return filename

# Generate RTF File with random image, and URL
def generate_rtf(url_type, include_image, include_sensitive_link):
    filename = f"generated_rtf_file_{random.randint(1000, 9999)}.rtf"
    file_path = os.path.join(FILE_STORAGE, filename)

    random_image_path = generate_image("jpg", url_type, include_sensitive_link) if include_image else None

    with open(file_path, "w") as f:
        # Start RTF document with basic content structure
        f.write("{\\rtf1\\ansi\\deff0 {\\fonttbl {\\f0 Courier;}}\n")
        f.write("\\fs24 Random RTF File Content\\par\n")

        # Add URL if provided
        if url_type != 'none':
            random_url = generate_random_url(url_type)
            f.write(f"Here is a URL: \\ul {random_url} \\ulnone\\par\n")  # Underlined URL

        # Add sensitive link if specified
        if include_sensitive_link == "on":
            sensitive_url = r"\\\\server\\share"  # Raw string to avoid issues with backslashes
            f.write(f"Here is a server share: \\ul {sensitive_url} \\ulnone\\par\n")  # Underlined Sensitive Link

        # Add image path if specified
        if random_image_path == "on":
            f.write(f"Random Image Path: app/data/generated_files/generated_file_1751.jpg \\par\n")

        # End RTF content
        f.write("}\n")  # Close RTF content

    
    return filename

# Generate Excel Spreadsheet (.xlsx) with a random URL
def generate_xlsx(url_type, include_image, include_sensitive_link):
    filename = f"generated_xlsx_file_{random.randint(1000, 9999)}.xlsx"
    file_path = os.path.join(FILE_STORAGE, filename)

    wb = Workbook()
    ws = wb.active
    ws.append(["This is an Excel file."])


    # Include include_sensitive_link if specified
    if include_sensitive_link == "on":
        sensitive_link = "\\server\share"
        logging.info(f" The sensitive Link is: {sensitive_link}")
        # Add the hyperlink to the Excel cell
        ws['A4'] = "Click here to visit the link"
        ws['A4'].hyperlink = sensitive_link  # This makes the text a clickable hyperlink
        ws['A4'].style = 'Hyperlink'  # Optional: applies the built-in hyperlink style
        
    # Generate a random URL based on the selected type
    if url_type != 'none':
        random_url = generate_random_url(url_type)
        logging.info(f" The generated Link is: {random_url}")
        # Add the hyperlink to the Excel cell
        ws['A2'] = "Click here to visit the link"
        ws['A2'].hyperlink = random_url  # This makes the text a clickable hyperlink
        ws['A2'].style = 'Hyperlink'  # Optional: applies the built-in hyperlink style

    # Adding an image if selected
    if include_image == "on":
        image_name = generate_image("jpg", url_type, include_sensitive_link)
        random_image_path = os.path.join(FILE_STORAGE, image_name)
        img = xlsImage(random_image_path)  # openpyxl Image object
        img.anchor = 'A10'  # Position the image in the cell (A1)
        ws.add_image(img)

    wb.save(file_path)

    return filename

def generate_image(file_type, url_type, include_sensitive_link):
    filename = f"generated_{file_type}_file_{random.randint(1000, 9999)}.{file_type}"
    file_path = os.path.join(FILE_STORAGE, filename)

    # Ensure the image is created in RGB mode (required for JPG)
    img = Image.new("RGB", (500, 500), (255, 255, 255))  # RGB mode
    draw = ImageDraw.Draw(img)

    # Add a gradient background (left to right)
    for i in range(500):
        r = int(255 * i / 500)
        g = int(255 * (500 - i) / 500)
        b = 100  # constant blue value for gradient
        draw.line((i, 0, i, 500), fill=(r, g, b))

    # Add some random shapes for complexity
    for _ in range(10):  # Add 10 random shapes
        shape_type = random.choice(["ellipse", "rectangle", "polygon"])
        x1, y1 = random.randint(0, 400), random.randint(0, 400)
        x2, y2 = random.randint(x1 + 20, 500), random.randint(y1 + 20, 500)
        
        if shape_type == "ellipse":
            draw.ellipse([x1, y1, x2, y2], outline="black", width=3)
        elif shape_type == "rectangle":
            draw.rectangle([x1, y1, x2, y2], outline="blue", width=3)
        elif shape_type == "polygon":
            draw.polygon([x1, y1, x2, y2, random.randint(0, 500), random.randint(0, 500)], outline="green", width=3)

    # Add dynamic text with varied size, fonts, and color
    font = ImageFont.load_default()  # You can specify a TTF font file here for a better font
    text = f"Random {file_type.upper()} Image"
    font_size = random.randint(20, 50)
    
    try:
        font = ImageFont.truetype("arial.ttf", font_size)  # Try loading a nicer font
    except IOError:
        font = ImageFont.load_default()  # Fallback if the font isn't available

    # Use textbbox to get the bounding box of the text
    text_bbox = draw.textbbox((0, 0), text, font=font)
    text_width = text_bbox[2] - text_bbox[0]
    text_height = text_bbox[3] - text_bbox[1]

    x = (500 - text_width) // 2
    y = (500 - text_height) // 4
    text_color = (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))
    draw.text((x, y), text, fill=text_color, font=font)
                
    # Add the URL text at the bottom
    if url_type != 'none':
        random_url = generate_random_url(url_type)
        url_text = f"{random_url}"
        url_font = ImageFont.load_default()
        url_text_bbox = draw.textbbox((0, 0), url_text, font=url_font)
        url_text_width = url_text_bbox[2] - url_text_bbox[0]
        draw.text(((500 - url_text_width) // 2, 450), url_text, font=url_font, fill=(0, 0, 0))

    # Include include_sensitive_link if specified
    if include_sensitive_link == "on":
        sensitive_link = "\\server\share"
        link_text = f"{sensitive_link}"
        link_font = ImageFont.load_default()
        link_text_bbox = draw.textbbox((0, 0), link_text, font=link_font)
        link_text_width = link_text_bbox[2] - link_text_bbox[0]
        draw.text(((500 - link_text_width) // 2, 450), link_text, font=link_font, fill=(0, 0, 0))
        
    # Save the image in the requested format (JPG or others)
    if file_type.lower() == 'jpg':
        img.save(file_path, format='JPEG')
    else:
        img.save(file_path, format=file_type.upper())  # Use the appropriate format for other types

    return filename

def generate_dylib(url_type=None):
    # Generate a unique filename for the dylib
    filename = f"generated_dylib_file_{random.randint(1000, 9999)}.dylib"
    c_file_path = os.path.join(FILE_STORAGE, f"{filename}.c")
    dylib_file_path = os.path.join(FILE_STORAGE, filename)
    # Create a random string or value to make the EXE unique
    random_value = random.randint(1000, 9999)
    
    # Simple C code to show a message
    c_code = f"""
#include <stdio.h>
#include <stdlib.h>

int main() {{
    // Generate a unique message with a random number
    int unique_value = {random_value};
    printf("This is a unique dylib message! The random value is: %d\\n", unique_value);
    
    // Wait for user to press Enter before closing
    printf("Press Enter to exit...\\n");
    getchar();  // Wait for the user to press Enter

    return 0;
}}
"""

    # Write the C code to a file
    with open(c_file_path, "w") as c_file:
        c_file.write(c_code)

    # Use GCC to compile the C file into an EXE
    try:
        subprocess.run(['gcc', c_file_path, '-o', dylib_file_path], check=True)
    except subprocess.CalledProcessError as e:
        logging.info(f"Error during dylib generation: {e}")
        return None

    # Clean up the C file after compilation
    os.remove(c_file_path)

    # Return the generated EXE filename
    return filename if os.path.exists(dylib_file_path) else None

# Generate Windows EXE - sudo apt-get install mingw-w64 on Kali or Ubuntu
def generate_exe(url_type=None):
    # Generate a unique filename for the EXE
    filename = f"generated_exe_file_{random.randint(1000, 9999)}.exe"
    c_file_path = os.path.join(FILE_STORAGE, f"{filename}.c")
    exe_file_path = os.path.join(FILE_STORAGE, filename)

    # Create a random string or value to make the EXE unique
    random_value = random.randint(1000, 9999)
    
    # Simple C code to show a message and include a unique value
    c_code = f"""
#include <stdio.h>
#include <stdlib.h>

int main() {{
    // Generate a unique message with a random number
    int unique_value = {random_value};
    printf("This is a unique EXE message! The random value is: %d\\n", unique_value);
    
    // Wait for user to press Enter before closing
    printf("Press Enter to exit...\\n");
    getchar();  // Wait for the user to press Enter

    return 0;
}}
"""

    # Write the C code to a file
    with open(c_file_path, "w") as c_file:
        c_file.write(c_code)

    # Use MinGW-w64 to compile the C file into an EXE (cross-compilation for Windows)
    try:
        subprocess.run(['x86_64-w64-mingw32-gcc', c_file_path, '-o', exe_file_path], check=True)
    except subprocess.CalledProcessError as e:
        logging.info(f"Error during EXE generation: {e}")
        return None

    # Clean up the C file after compilation
    os.remove(c_file_path)

    # Return the generated EXE filename
    return filename if os.path.exists(exe_file_path) else None

def generate_elf(url_type=None):
        filename = f"generated_elf_file_{random.randint(1000, 9999)}.elf"
        c_file_path = os.path.join(FILE_STORAGE, f"{filename}.c")
        elf_file_path = os.path.join(FILE_STORAGE, filename)

        random_value = random.randint(1000, 9999)

        # C code for the shared object
        c_code = f"""
        #include <stdio.h>

        void display_message() {{
            int unique_value = {random_value};
           printf("This is a unique ELF message! The random value is: %d\\n", unique_value);
        }}
        """

        with open(c_file_path, "w") as c_file:
            c_file.write(c_code)

        # Compile the C file into a shared object
        try:
            subprocess.run(['gcc', '-shared', '-fPIC', c_file_path, '-o', elf_file_path],
                        check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        except subprocess.CalledProcessError as e:
            logging.info(f"Error during ELF generation: {e}")
            return None
        
        os.remove(c_file_path)

        # Return the generated EXE filename
        return filename if os.path.exists(elf_file_path) else None

